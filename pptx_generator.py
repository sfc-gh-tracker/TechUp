import os
from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = os.path.dirname(os.path.abspath(__file__))

slides_spec = [
    {
        "folder": "Adaptive Right-Sizing",
        "title": "Adaptive Warehouse Right-Sizing",
        "bullets": [
            "Goal: auto-tune warehouse size by observed load patterns",
            "Inputs: staged WAREHOUSE_METERING (credits_used by hour)",
            "Policy DT: RIGHT_SIZING_POLICY_DT → per-warehouse/hour recommendation",
            "Executor: APPLY_RIGHT_SIZING() applies size + optional multi-cluster"
        ],
        "details_title": "How it works",
        "details": [
            "Ingestion: Task merges ACCOUNT_USAGE.WAREHOUSE_METERING → TECHUP.AUDIT.WAREHOUSE_METERING_STG (change tracking)",
            "Signal: Aggregate credits_used into hourly buckets per warehouse",
            "Policy logic: map avg(credits_used) ranges → SMALL/MEDIUM/LARGE sizing and multi-cluster toggle",
            "Governance: all changes logged in RIGHT_SIZING_LOG with status, DDL, error",
            "Orchestration: APPLY_RIGHT_SIZING_TASK executes on-the-hour against current hour recommendation",
            "Safety: thresholds are conservative; tune sizing cutoffs per environment; dry-run by commenting execute immediate"
        ],
    },
    {
        "folder": "Pipeline Factory",
        "title": "Natural Language to SQL Pipeline Factory",
        "bullets": [
            "Streamlit app uses Cortex to generate validated SQL",
            "Writes SQL into PIPELINE_CONFIG (db-level target DT)",
            "RUN_PIPELINE_FACTORY creates/refreshes Dynamic Tables",
            "Searchable DB/Schema, allowed tables multi-select"
        ],
        "details_title": "How it works",
        "details": [
            "Discovers schema via INFORMATION_SCHEMA with token-safe truncation",
            "Cortex generates SQL; validation fixes joins/columns; preview 3 rows",
            "Insert snippet with target_dt_database + target_dt_name",
            "RUN_PIPELINE_FACTORY builds DT: create or replace dynamic table <db>..<name>"
        ],
    },
    {
        "folder": "Query Pattern Optimizer",
        "title": "Query Pattern Optimizer",
        "bullets": [
            "Stages QUERY_HISTORY into TECHUP.QPO_AUDIT.QUERY_HISTORY_STG",
            "Aggregates patterns and flags heavy scans/spillage",
            "Recommendations DT emits suggested DDL actions",
            "Executor procedure applies reviewed actions via task"
        ],
        "details_title": "How it works",
        "details": [
            "Task merges ACCOUNT_USAGE.QUERY_HISTORY → QUERY_HISTORY_STG (DT-safe)",
            "QPO_USAGE aggregation DT extracts tables, bytes_scanned, patterns",
            "Recommendations DT proposes clustering or warehouse sizing",
            "Pending DDL DT + QPO_RUN_ACTIONS() execute reviewed DDL (logs)"
        ],
    },
    {
        "folder": "Performance Monitor",
        "title": "Self-Optimizing Performance Monitor",
        "bullets": [
            "Stages WAREHOUSE_METERING into TECHUP.AUDIT.WAREHOUSE_METERING_STG",
            "WAREHOUSE_PERFORMANCE_DT joins query + metering signals",
            "OPTIMIZATION_RECOMMENDATIONS_DT proposes improvements",
            "PENDING_DDL_ACTIONS_DT + RUN_PENDING_ACTIONS() for automation"
        ],
        "details_title": "How it works",
        "details": [
            "Ingestion task maintains metering staging with change tracking",
            "Performance DT correlates bytes_scanned/spillage with time/warehouse",
            "Recommendations DT flags high spillage/scans/utilization (credits_used)",
            "Executor procedure applies actions; ACTION_LOG records results"
        ],
    },
]


def add_bullets_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(18)

def create_pptx_with_details(title: str, bullets: list[str], details_title: str | None, details: list[str] | None, out_path: str) -> None:
    prs = Presentation()
    add_bullets_slide(prs, title, bullets)
    if details and details_title:
        add_bullets_slide(prs, details_title, details)
    prs.save(out_path)


def main() -> None:
    for spec in slides_spec:
        folder = os.path.join(ROOT, spec["folder"])
        os.makedirs(folder, exist_ok=True)
        out_path = os.path.join(folder, "overview.pptx")
        create_pptx_with_details(
            spec["title"],
            spec["bullets"],
            spec.get("details_title"),
            spec.get("details"),
            out_path,
        )
        print(f"Wrote {out_path}")


if __name__ == "__main__":
    main()


